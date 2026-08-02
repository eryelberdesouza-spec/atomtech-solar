"""
Preenche templates/template_relatorio_gestao_energia.pptx com um dicionário
{TOKEN: valor} e salva o relatório final do mês/cliente.

Uso:
    python3 src/fill_template.py dados.json saida.pptx
"""
import sys
import json
from pathlib import Path
from pptx import Presentation

TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "templates" / "template_relatorio_gestao_energia.pptx"


def fill_text_frame(tf, data):
    for p in tf.paragraphs:
        for r in p.runs:
            for token, value in data.items():
                placeholder = "{{" + token + "}}"
                if placeholder in r.text:
                    r.text = r.text.replace(placeholder, str(value))


def walk_shapes(shapes, data):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            walk_shapes(shape.shapes, data)
            continue
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    fill_text_frame(cell.text_frame, data)
            continue
        if shape.has_text_frame:
            fill_text_frame(shape.text_frame, data)


def fill(template_path, data, out_path):
    prs = Presentation(template_path)
    for slide in prs.slides:
        walk_shapes(slide.shapes, data)
    prs.save(out_path)


if __name__ == "__main__":
    data_path, out_path = sys.argv[1], sys.argv[2]
    with open(data_path, encoding="utf-8") as f:
        data = json.load(f)
    fill(TEMPLATE_PATH, data, out_path)
    print("Salvo:", out_path)
